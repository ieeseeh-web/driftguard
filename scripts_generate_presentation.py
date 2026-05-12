from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path("agent/DriftGuard_Architecture_Presentation.pptx")

# MiniMax-inspired design tokens from DESIGN.md
COLORS = {
    "primary": "0B0B0D",
    "canvas": "FFFFFF",
    "surface": "F7F7F5",
    "surface_soft": "FAFAF8",
    "hairline": "E7E7E4",
    "ink": "111114",
    "charcoal": "34343A",
    "steel": "6F727A",
    "stone": "9A9A9A",
    "coral": "FF5A3D",
    "magenta": "D92D8A",
    "blue": "2563FF",
    "blue_deep": "1848D6",
    "blue_200": "EAF1FF",
    "purple": "6D4DFF",
    "success_bg": "E9F8EE",
    "success_text": "137A3D",
    "error": "D45656",
}

FONT = "DM Sans"
WIDE = (13.333, 7.5)


def rgb(key: str) -> RGBColor:
    h = COLORS[key]
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color="canvas"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=18, weight=False, color="ink", align=PP_ALIGN.LEFT, line_spacing=1.1):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = weight
    run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        add_pill(slide, kicker, 0.72, 0.42, 1.8, 0.34, fill="primary", text_color="canvas", size=9)
    add_text(slide, title, 0.7, 0.82, 10.8, 0.9, size=30, weight=True, color="ink")
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.62, 9.5, 0.46, size=13, color="steel")
    add_slide_number(slide)


def add_slide_number(slide):
    idx = len(slide.part.package.presentation_part.presentation.slides)
    add_text(slide, f"{idx:02d}", 12.35, 7.06, 0.35, 0.2, size=8, color="stone", align=PP_ALIGN.RIGHT)


def add_pill(slide, text, x, y, w, h, fill="canvas", text_color="ink", line="hairline", size=10, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line if line else fill)
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(text_color)
    return shape


def add_card(slide, x, y, w, h, fill="surface", line="hairline", radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def add_bullets(slide, items, x, y, w, h, size=15, color="charcoal", bullet_color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(8)
        p.line_spacing = 1.12
    return box


def add_table(slide, headers, rows, x, y, w, h, col_widths=None, font_size=11):
    nrows = len(rows) + 1
    ncols = len(headers)
    table_shape = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb("surface")
        cell.text = head
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT; p.font.bold = True; p.font.size = Pt(font_size); p.font.color.rgb = rgb("steel")
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb("canvas")
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT; p.font.size = Pt(font_size); p.font.color.rgb = rgb("charcoal")
    return table_shape


def connector(slide, x1, y1, x2, y2, color="steel"):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(1.2)
    return line


def node(slide, text, x, y, w, h, fill="canvas", color="ink", line="hairline", size=11):
    add_card(slide, x, y, w, h, fill=fill, line=line)
    return add_text(slide, text, x+0.12, y+0.14, w-0.24, h-0.2, size=size, weight=True, color=color, align=PP_ALIGN.CENTER)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_pill(slide, "DRIFTGUARD", 0.72, 0.54, 1.65, 0.36, fill="primary", text_color="canvas", size=9)
    add_text(slide, "Agent Drift\nJudge & Agent\nArchitecture", 0.72, 1.25, 7.0, 2.45, size=44, weight=True, color="ink", line_spacing=0.9)
    add_text(slide, "시스템 가드레일 방식과 AI 평가 에이전트 방식 비교", 0.78, 4.02, 6.6, 0.42, size=15, color="steel")
    # Four vibrant product-style cards
    cards = [("System", "Guardrail", "coral"), ("Agent", "Reviewer", "magenta"), ("Drift", "Score", "blue"), ("Audit", "Log", "purple")]
    for i, (a, b, c) in enumerate(cards):
        x = 7.95 + (i % 2) * 2.25; y = 1.1 + (i // 2) * 2.15
        add_card(slide, x, y, 1.95, 1.82, fill=c, line=c)
        add_text(slide, a, x+0.22, y+0.36, 1.45, 0.36, size=18, weight=True, color="canvas")
        add_text(slide, b, x+0.22, y+0.86, 1.45, 0.28, size=10, color="canvas")
    add_text(slide, "MiniMax-inspired deck style · DM Sans · pill CTAs · vibrant product cards", 0.78, 6.83, 9.4, 0.22, size=9, color="stone")
    add_slide_number(slide)


def problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Agent Drift란 무엇인가", "초기 사용자 목표·역할·제약·정책·맥락에서 점진적으로 벗어나는 현상", "Problem")
    add_card(slide, 0.75, 2.25, 3.2, 3.25, fill="surface")
    add_text(slide, "핵심 위험", 1.0, 2.52, 1.6, 0.28, size=16, weight=True)
    add_bullets(slide, ["장기 작업에서 작은 오해가 누적", "도구 호출로 외부 상태 변경", "메모리 저장으로 오류가 지속", "다중 에이전트 handoff에서 목표 변형"], 1.0, 3.0, 2.45, 1.95, size=12)
    # Flow visual
    xs = [4.75, 6.35, 7.95, 9.55]
    labels = ["Original\nGoal", "Plan", "Tool /\nMemory", "Drifted\nOutcome"]
    colors = ["canvas", "surface", "blue_200", "coral"]
    for x, label, c in zip(xs, labels, colors):
        node(slide, label, x, 3.0, 1.15, 0.8, fill=c, color="canvas" if c == "coral" else "ink")
    for i in range(3):
        connector(slide, xs[i]+1.15, 3.4, xs[i+1], 3.4)
    add_pill(slide, "scope expansion", 6.2, 4.23, 1.7, 0.35, fill="magenta", text_color="canvas", line="magenta", size=8)
    add_pill(slide, "unsafe action", 8.0, 4.23, 1.55, 0.35, fill="purple", text_color="canvas", line="purple", size=8)
    add_text(slide, "예: README만 수정 요청 → architecture.md까지 변경 · 삭제 금지 → rm -rf 실행", 4.75, 5.2, 6.0, 0.34, size=12, color="steel")


def drift_types_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Drift 유형", "Drift는 단일 응답 문제가 아니라 실행 과정 전체에서 발생한다", "Taxonomy")
    headers = ["유형", "설명", "예시"]
    rows = [
        ["Goal", "원래 목표에서 벗어남", "요청 범위 외 파일 수정"],
        ["Instruction", "명시 지시 누락", "삭제 금지 무시"],
        ["Tool", "위험/불필요 도구 사용", "승인 없이 배포"],
        ["Memory", "부적절한 기억 저장", "일시적 선호를 영구 저장"],
        ["Multi-Agent", "전달 과정 목표 변형", "Planner 지시가 원본과 달라짐"],
        ["Safety", "승인/민감정보/외부 영향", "토큰 로그 저장"],
    ]
    add_table(slide, headers, rows, 0.85, 2.15, 11.6, 3.8, col_widths=[2.0, 4.25, 5.35], font_size=10)
    add_pill(slide, "Goal", 1.0, 6.34, 0.9, 0.32, fill="coral", text_color="canvas", line="coral", size=8)
    add_pill(slide, "Tool", 2.05, 6.34, 0.9, 0.32, fill="blue", text_color="canvas", line="blue", size=8)
    add_pill(slide, "Memory", 3.1, 6.34, 1.15, 0.32, fill="purple", text_color="canvas", line="purple", size=8)
    add_pill(slide, "Safety", 4.4, 6.34, 1.0, 0.32, fill="primary", text_color="canvas", line="primary", size=8)


def system_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "접근 방식 1: 시스템 아키텍처", "Agent Runtime에 Judge Layer, Policy Engine, Tool Guard, Memory Guard를 삽입", "System")
    node(slide, "User\nRequest", 0.8, 3.1, 1.25, 0.75, fill="surface")
    node(slide, "Agent\nRuntime", 2.6, 2.65, 1.6, 1.65, fill="primary", color="canvas", line="primary")
    node(slide, "Judge\nLayer", 4.95, 1.85, 1.45, 0.78, fill="blue", color="canvas", line="blue")
    node(slide, "Policy\nEngine", 4.95, 3.0, 1.45, 0.78, fill="coral", color="canvas", line="coral")
    node(slide, "Tool\nGuard", 4.95, 4.15, 1.45, 0.78, fill="magenta", color="canvas", line="magenta")
    node(slide, "Memory\nGuard", 6.95, 4.15, 1.45, 0.78, fill="purple", color="canvas", line="purple")
    node(slide, "Continue", 8.95, 2.0, 1.3, 0.55, fill="success_bg", color="success_text")
    node(slide, "Ask User", 8.95, 3.05, 1.3, 0.55, fill="blue_200", color="blue_deep")
    node(slide, "Stop &\nAudit", 8.95, 4.1, 1.3, 0.72, fill="primary", color="canvas", line="primary")
    for args in [(2.05,3.48,2.6,3.48),(4.2,3.48,4.95,3.38),(6.4,3.38,8.95,3.32),(6.4,4.5,6.95,4.5),(8.4,4.5,8.95,4.48)]: connector(slide,*args)
    add_card(slide, 10.7, 1.85, 1.55, 3.0, fill="surface")
    add_text(slide, "강점", 10.95, 2.12, 0.8, 0.24, size=14, weight=True)
    add_bullets(slide, ["자동 보호", "운영 통합", "실시간 차단"], 10.95, 2.6, 1.05, 0.9, size=10)
    add_text(slide, "제약", 10.95, 3.9, 0.8, 0.24, size=14, weight=True)
    add_bullets(slide, ["통합 비용", "런타임 의존"], 10.95, 4.35, 1.05, 0.6, size=10)


def components_slide(prs, title, rows, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, title, None, kicker)
    add_table(slide, ["컴포넌트", "역할"], rows, 0.85, 1.82, 11.5, 4.35, col_widths=[3.1, 8.4], font_size=10)
    add_text(slide, "flat-with-borders documentation surface · compact rows · 16px white cards", 0.9, 6.48, 7.5, 0.25, size=9, color="stone")


def flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "시스템 방식 실행 흐름", "최종 응답에서만 평가하면 늦다 — 계획, 도구 호출, 메모리 업데이트 시점에 체크한다", "Flow")
    steps = ["User\nRequest", "Agent\nPlanning", "Judge\nEvaluation", "Policy\nDecision", "Agent\nExecution", "Tool / Memory\nGuard", "Final\nEvaluation", "Continue /\nAsk / Stop"]
    x0, y = 0.7, 3.0
    for i, s in enumerate(steps):
        x = x0 + i*1.55
        fill = "primary" if i in [2,3,6] else ("coral" if i == 7 else "surface")
        color = "canvas" if fill in ["primary", "coral"] else "ink"
        node(slide, s, x, y, 1.15, 0.78, fill=fill, color=color, line=fill if fill != "surface" else "hairline", size=9)
        if i < len(steps)-1:
            connector(slide, x+1.15, y+0.39, x+1.55, y+0.39)
    add_pill(slide, "checkpoint-based evaluation", 4.65, 4.35, 2.6, 0.36, fill="blue_200", text_color="blue_deep", line="blue_200", size=9)


def agent_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "접근 방식 2: 에이전트 아키텍처", "DriftGuard를 독립적인 평가 에이전트로 실행하고, 설명과 수정 가이드를 제공", "Agent")
    node(slide, "Review\nRequest", 0.85, 3.1, 1.35, 0.75, fill="surface")
    node(slide, "Input\nNormalizer", 2.65, 3.1, 1.45, 0.75, fill="blue_200", color="blue_deep", line="blue_200")
    node(slide, "Rule-based\nEvaluator", 4.55, 2.2, 1.55, 0.75, fill="surface")
    node(slide, "DriftGuard\nAgent", 4.55, 3.55, 1.55, 0.85, fill="primary", color="canvas", line="primary")
    node(slide, "Drift Type\nClassifier", 6.65, 2.25, 1.55, 0.75, fill="coral", color="canvas", line="coral")
    node(slide, "Guidance\nGenerator", 6.65, 3.65, 1.55, 0.75, fill="magenta", color="canvas", line="magenta")
    node(slide, "Markdown\nReport", 9.0, 2.1, 1.35, 0.72, fill="canvas")
    node(slide, "Structured\nJSON", 9.0, 3.1, 1.35, 0.72, fill="canvas")
    node(slide, "JSONL\nAudit Log", 9.0, 4.1, 1.35, 0.72, fill="canvas")
    for args in [(2.2,3.48,2.65,3.48),(4.1,3.48,4.55,3.9),(4.1,3.48,4.55,2.58),(6.1,3.9,6.65,4.02),(6.1,3.9,6.65,2.62),(8.2,2.62,9.0,2.46),(8.2,3.95,9.0,3.46),(8.2,3.95,9.0,4.46)]: connector(slide,*args)
    add_card(slide, 10.85, 2.0, 1.35, 2.8, fill="surface")
    add_text(slide, "핵심", 11.1, 2.25, 0.8, 0.22, size=13, weight=True)
    add_bullets(slide, ["리뷰어", "코치", "감사자"], 11.1, 2.75, 0.9, 0.8, size=10)


def example_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "에이전트 방식 실행 예시", "실행 로그를 review-agent 입력으로 넣으면 Drift 유형·점수·권고가 나온다", "Demo")
    add_card(slide, 0.85, 1.85, 5.55, 3.95, fill="primary", line="primary")
    add_text(slide, "CLI", 1.15, 2.15, 0.5, 0.22, size=12, weight=True, color="canvas")
    code = "PYTHONPATH=backend/src python3 -m driftguard.cli review-agent \\\n  --input backend/examples/agent-review-execution-log.json \\\n  --log backend/logs/agent-reviews.jsonl"
    add_text(slide, code, 1.15, 2.62, 4.8, 0.8, size=10, color="canvas")
    add_text(slide, "Input", 1.15, 3.85, 0.7, 0.22, size=12, weight=True, color="canvas")
    add_text(slide, '{"execution_log": ["Ran command: rm -rf docs/old"]}', 1.15, 4.32, 4.8, 0.34, size=10, color="canvas")
    add_card(slide, 7.05, 1.85, 4.95, 3.95, fill="surface")
    add_text(slide, "Output", 7.35, 2.15, 0.9, 0.22, size=12, weight=True)
    rows = [["drift_types", "tool, safety"], ["score", "0.88"], ["risk", "critical"], ["recommendation", "stop"]]
    add_table(slide, ["Field", "Value"], rows, 7.35, 2.62, 4.15, 2.0, col_widths=[1.55, 2.6], font_size=10)
    add_pill(slide, "삭제 작업은 중단하고 사용자 확인을 받으세요", 7.35, 5.02, 3.75, 0.36, fill="coral", text_color="canvas", line="coral", size=9)


def comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "시스템 방식 vs 에이전트 방식", "브레이크와 리뷰어/코치 — 서로 다른 목적의 두 레이어", "Compare")
    rows = [
        ["목적", "런타임 보호", "리뷰/진단/가이드"],
        ["위치", "Agent Runtime 내부/인접", "독립 Agent 또는 sub-agent"],
        ["실행 시점", "실시간/중간 단계", "사전·중간 리뷰, 사후 감사"],
        ["출력", "승인/차단/재계획", "리포트, 설명, 수정 가이드"],
        ["장점", "자동 보호, 운영 통합", "설명 가능성, 점진적 도입"],
        ["단점", "통합 비용", "실시간 차단력 제한"],
    ]
    add_table(slide, ["항목", "시스템 방식", "에이전트 방식"], rows, 0.75, 1.8, 11.85, 4.35, col_widths=[2.0, 4.85, 5.0], font_size=9)
    add_pill(slide, "System = Brake", 1.0, 6.45, 1.8, 0.36, fill="primary", text_color="canvas", line="primary", size=9)
    add_pill(slide, "Agent = Reviewer / Coach", 3.05, 6.45, 2.55, 0.36, fill="coral", text_color="canvas", line="coral", size=9)


def strategy_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "두 방식의 결합 전략", "먼저 에이전트 방식으로 기준을 쌓고, 검증된 기준을 시스템 가드레일로 승격", "Strategy")
    phases = [("Phase 1", "Agent Review CLI", "coral"), ("Phase 2", "Sub-agent / CI Review", "magenta"), ("Phase 3", "Runtime Hook", "blue"), ("Phase 4", "Dashboard + Policy", "purple")]
    for i, (p, t, c) in enumerate(phases):
        x = 0.85 + i*3.05
        add_card(slide, x, 2.65, 2.35, 1.75, fill=c, line=c)
        add_text(slide, p, x+0.25, 2.98, 1.2, 0.24, size=12, weight=True, color="canvas")
        add_text(slide, t, x+0.25, 3.52, 1.8, 0.45, size=16, weight=True, color="canvas")
        if i < 3:
            connector(slide, x+2.35, 3.53, x+3.05, 3.53)
    add_text(slide, "점진적 도입: 개발 리뷰 → 자동 검증 → 런타임 보호 → 운영 지표화", 0.9, 5.35, 8.5, 0.36, size=15, color="charcoal")


def mvp_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "현재 MVP 구현 상태", "CLI, 샘플, 로그, 테스트까지 동작하는 최소 기능 제품", "MVP")
    rows = [["Rule-based evaluator", "완료"], ["review-agent CLI", "완료"], ["Markdown + JSON 출력", "완료"], ["JSONL 로그 저장", "완료"], ["Tool / Memory / Final 샘플", "완료"], ["Handoff / Execution Log 샘플", "완료"], ["테스트", "13개 통과"], ["LLM Judge 연동", "향후 과제"], ["Dashboard", "향후 과제"]]
    add_table(slide, ["항목", "상태"], rows, 0.85, 1.75, 7.0, 4.8, col_widths=[4.7, 2.3], font_size=10)
    add_card(slide, 8.4, 2.05, 3.45, 2.7, fill="success_bg", line="success_bg")
    add_text(slide, "Verification", 8.82, 2.45, 1.5, 0.28, size=15, weight=True, color="success_text")
    add_text(slide, "PYTHONPATH=backend/src python3 -m unittest discover -s tests -v\n\nRan 13 tests OK", 8.82, 3.05, 2.4, 1.0, size=12, color="success_text")


def demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Demo 시나리오", "세 가지 대표 Drift를 샘플 입력으로 재현하고 평가한다", "Demo")
    demos = [("Tool Drift", "삭제하지 말고 요약해줘 → rm -rf docs/old", "critical · stop", "coral"), ("Memory Drift", "오늘은 짧게 → 항상 짧은 답변 선호", "high · skip_memory", "purple"), ("Multi-Agent Drift", "README만 수정 → architecture.md도 정리", "multi_agent · stop", "blue")]
    for i, (title, body, result, c) in enumerate(demos):
        x = 0.85 + i*4.05
        add_card(slide, x, 2.2, 3.35, 3.25, fill="surface")
        add_pill(slide, title, x+0.28, 2.52, 1.6, 0.34, fill=c, text_color="canvas", line=c, size=9)
        add_text(slide, body, x+0.32, 3.15, 2.55, 0.75, size=13, weight=True)
        add_text(slide, result, x+0.32, 4.55, 2.1, 0.32, size=16, weight=True, color=c)


def roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_title(slide, "Roadmap", "Agent Review CLI에서 AgentOps 지표 추적까지", "Roadmap")
    items = ["Agent Review CLI 안정화", "JSON Schema 검증과 로그 마스킹", "LLM Judge 연동", "GitHub PR / CI 리뷰 모드", "OpenClaw sub-agent 실행", "AgentOps 대시보드와 Drift 지표 추적"]
    for i, item in enumerate(items):
        y = 1.75 + i*0.72
        add_pill(slide, str(i+1), 1.0, y, 0.45, 0.36, fill="primary", text_color="canvas", line="primary", size=9)
        add_text(slide, item, 1.7, y+0.04, 6.4, 0.25, size=15, weight=True if i < 3 else False)
        if i < len(items)-1:
            connector(slide, 1.225, y+0.36, 1.225, y+0.72)
    add_card(slide, 8.45, 2.15, 2.85, 2.85, fill="coral", line="coral")
    add_text(slide, "North Star", 8.85, 2.65, 1.7, 0.28, size=15, weight=True, color="canvas")
    add_text(slide, "Drift를 차단하는 것에서 끝나지 않고, 에이전트를 안전한 목표 경로로 되돌린다.", 8.85, 3.3, 1.95, 1.0, size=14, weight=True, color="canvas")


def conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
    add_pill(slide, "Conclusion", 0.75, 0.55, 1.35, 0.34, fill="primary", text_color="canvas", size=9)
    add_text(slide, "DriftGuard의 목표는\n에이전트를 멈추게 하는 것이 아니라,\n원래 목표와 안전한 경로로 되돌리는 것이다.", 0.8, 1.32, 8.9, 2.25, size=34, weight=True, line_spacing=0.95)
    add_card(slide, 0.9, 4.35, 3.1, 1.35, fill="surface")
    add_text(slide, "System Architecture", 1.2, 4.68, 2.0, 0.25, size=15, weight=True)
    add_text(slide, "런타임 보호에 강함", 1.2, 5.12, 2.0, 0.22, size=12, color="steel")
    add_card(slide, 4.4, 4.35, 3.1, 1.35, fill="surface")
    add_text(slide, "Agent Architecture", 4.7, 4.68, 2.0, 0.25, size=15, weight=True)
    add_text(slide, "설명 가능한 리뷰에 강함", 4.7, 5.12, 2.0, 0.22, size=12, color="steel")
    add_card(slide, 7.9, 4.35, 3.1, 1.35, fill="primary", line="primary")
    add_text(slide, "Combined", 8.2, 4.68, 2.0, 0.25, size=15, weight=True, color="canvas")
    add_text(slide, "평가 · 완화 · 감사", 8.2, 5.12, 2.0, 0.22, size=12, color="canvas")
    add_slide_number(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0]); prs.slide_height = Inches(WIDE[1])
    # remove default slide? Presentation starts empty enough.
    title_slide(prs)
    problem_slide(prs)
    drift_types_slide(prs)
    system_overview_slide(prs)
    components_slide(prs, "시스템 아키텍처 컴포넌트", [["Agent Runtime", "계획, 실행, 도구 호출, 최종 응답 생성"], ["Judge Layer", "목표/역할/지시/도구/메모리 평가"], ["Policy Engine", "점수와 정책 기반으로 다음 행동 결정"], ["Tool Guard", "고위험 도구 호출 전 검증"], ["Memory Guard", "장기 메모리 저장 전 검증"], ["Evaluation Log", "감사와 운영 지표 기록"]], "Components")
    flow_slide(prs)
    agent_overview_slide(prs)
    components_slide(prs, "에이전트 아키텍처 컴포넌트", [["Agent Review Request", "평가 입력 JSON"], ["Input Normalizer", "로그/응답/도구 호출을 공통 형식으로 정리"], ["Rule-based Evaluator", "빠른 위험 신호 계산"], ["DriftGuard Agent", "Drift 진단, 근거, 수정 가이드 생성"], ["Guidance Generator", "실행 가능한 수정 방향 생성"], ["Output Layer", "Markdown, JSON, JSONL 로그 출력"]], "Components")
    example_slide(prs)
    comparison_slide(prs)
    strategy_slide(prs)
    mvp_slide(prs)
    demo_slide(prs)
    roadmap_slide(prs)
    conclusion_slide(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
