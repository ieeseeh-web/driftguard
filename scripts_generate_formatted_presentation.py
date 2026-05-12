from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt

from scripts_generate_deep_presentation import slides as CONTENT_SLIDES

OUT = Path("agent/DriftGuard_Agent_Drift_DeepDive_Presentation_formatted.pptx")

FONT = "Malgun Gothic"  # 맑은고딕
RED = RGBColor(234, 0, 44)
ORANGE = RGBColor(255, 122, 0)
GRAY25 = RGBColor(191, 191, 191)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(45, 45, 45)
LIGHT_GRAY = RGBColor(245, 245, 245)

SLIDE_W = Cm(33.867)  # 16:9 widescreen
SLIDE_H = Cm(19.05)
HEADER_H = Cm(1.58)
TITLE_X = Cm(0.9)
TITLE_Y = Cm(0.57)
TITLE_W = Cm(27.77)
TITLE_H = Cm(0.92)
BODY_X = Cm(0.9)
BODY_Y = Cm(1.95)
BODY_W = Cm(31.8)
BODY_H = Cm(16.45)


def set_white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_header(slide, title: str, idx: int):
    # Header area: white background, gray bottom rule, red/orange accent bars.
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.fill.background()

    red_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(0.22), HEADER_H)
    red_bar.fill.solid(); red_bar.fill.fore_color.rgb = RED; red_bar.line.fill.background()
    orange_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.22), 0, Cm(0.10), HEADER_H)
    orange_bar.fill.solid(); orange_bar.fill.fore_color.rgb = ORANGE; orange_bar.line.fill.background()

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, HEADER_H - Pt(1), SLIDE_W, Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = GRAY25; line.line.fill.background()

    box = slide.shapes.add_textbox(TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    tf = box.text_frame
    tf.clear(); tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = BLACK

    num = slide.shapes.add_textbox(Cm(31.5), Cm(0.58), Cm(1.5), Cm(0.5))
    nt = num.text_frame
    nt.clear(); nt.margin_left = 0; nt.margin_right = 0; nt.margin_top = 0; nt.margin_bottom = 0
    p = nt.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{idx:02d}"; r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = GRAY25


def add_l1(slide, text: str, x, y, w, h=Cm(0.65), color=BLACK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear(); tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.name = FONT; run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = color
    return box


def add_l2_bullets(slide, items: list[str], x, y, w, h, color=DARK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.line_spacing = 1.12
    return box


def add_text(slide, value: str, x, y, w, h, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, fill=LIGHT_GRAY, line=GRAY25):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.08
    return shape


def add_color_label(slide, text: str, x, y, w, h, fill=RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background(); shape.adjustments[0] = 0.5
    tf = shape.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE
    return shape


def add_code(slide, code: str, x, y, w, h):
    add_panel(slide, x, y, w, h, fill=RGBColor(250, 250, 250), line=GRAY25)
    # Keep body level-2 size per requested format.
    return add_text(slide, code, x + Cm(0.35), y + Cm(0.28), w - Cm(0.7), h - Cm(0.56), size=12, color=BLACK)


def add_table(slide, headers: list[str], rows: list[list[str]], x, y, w, h):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    tbl = shape.table
    for c, head in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = RED if c == 0 else ORANGE
        cell.text = head
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT; p.font.size = Pt(12); p.font.color.rgb = DARK
    return shape


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = GRAY25
    line.line.width = Pt(1.3)


def add_node(slide, label, x, y, w, h, fill=LIGHT_GRAY, text_color=BLACK):
    add_panel(slide, x, y, w, h, fill=fill, line=GRAY25)
    add_text(slide, label, x + Cm(0.15), y + Cm(0.18), w - Cm(0.3), h - Cm(0.3), size=12, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def add_flow(slide, labels: list[str], y=Cm(8.0)):
    n = len(labels)
    gap = Cm(0.25)
    w = (BODY_W - gap * (n - 1)) / n
    x = BODY_X
    colors = [WHITE, LIGHT_GRAY, RED, ORANGE, LIGHT_GRAY, WHITE, RED, ORANGE]
    for i, label in enumerate(labels):
        fill = colors[i % len(colors)]
        text_color = WHITE if fill in (RED, ORANGE) else BLACK
        add_node(slide, label, x, y, w, Cm(1.5), fill=fill, text_color=text_color)
        if i < n - 1:
            add_connector(slide, x + w, y + Cm(0.75), x + w + gap, y + Cm(0.75))
        x += w + gap


def add_simple_diagram(slide, kind: str):
    if kind == "system_arch":
        labels = ["User\nRequest", "Intent\nContract", "Agent\nRuntime", "Judge\nRouter", "Guards", "Policy\nEngine", "Action /\nAudit"]
        add_flow(slide, labels, y=Cm(7.0))
        return
    if kind == "agent_judge":
        labels = ["Agent\nTrace", "Observer", "Requirement\nExtractor", "Step\nEvaluator", "Evidence\nCollector", "Policy\nAdvisor", "Report"]
        add_flow(slide, labels, y=Cm(7.0))
        return
    if kind == "multi_judge":
        labels = ["Goal\nJudge", "Tool\nJudge", "Memory\nJudge", "Privacy\nJudge", "Aggregator", "Policy\nDecision"]
        add_flow(slide, labels, y=Cm(7.0))
        return
    labels = ["Input", "Goal", "Plan", "Tool", "Memory", "Handoff", "Action"]
    add_flow(slide, labels, y=Cm(7.0))


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    add_header(slide, "Agent Drift Detection & Judge Architecture", 1)
    add_l1(slide, "LLM-as-a-Judge와 Agent-as-a-Judge를 활용한 에이전트 목표 이탈 탐지·평가·완화 전략", BODY_X, Cm(3.0), BODY_W, Cm(1.0), color=BLACK)
    add_l2_bullets(slide, [
        "Agent Drift를 단순 응답 오류가 아닌 목표 유지·실행 경로 문제로 정의",
        "LLM-as-a-Judge와 Agent-as-a-Judge의 역할 차이 설명",
        "실제 구현에 필요한 Intent Contract, Guard, Policy, Audit 구조 제안",
    ], BODY_X, Cm(4.5), Cm(20), Cm(3.5))
    # Accent blocks using requested colors.
    x0 = Cm(22.0)
    for i, (label, fill) in enumerate([("Goal", RED), ("Tool", ORANGE), ("Memory", GRAY25)]):
        add_panel(slide, x0, Cm(3.0 + i * 2.4), Cm(7.5), Cm(1.7), fill=fill, line=fill)
        add_text(slide, label, x0 + Cm(0.35), Cm(3.35 + i * 2.4), Cm(5), Cm(0.6), size=16, bold=True, color=WHITE if fill != GRAY25 else BLACK)


def add_content_slide(prs, spec: dict, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    add_header(slide, spec["title"], idx)

    if spec.get("subtitle"):
        add_l1(slide, spec["subtitle"], BODY_X, BODY_Y, BODY_W, Cm(0.8), color=RED)
        top = Cm(3.0)
    else:
        top = BODY_Y

    if "diagram" in spec:
        add_simple_diagram(slide, spec["diagram"])
        if spec.get("bullets"):
            add_l1(slide, "주요 포인트", BODY_X, Cm(11.0), BODY_W, Cm(0.7))
            add_l2_bullets(slide, spec["bullets"][:4], BODY_X, Cm(11.9), BODY_W, Cm(3.2))
        return

    if "flow" in spec:
        add_flow(slide, spec["flow"], y=Cm(6.0))
        if spec.get("bullets"):
            add_l1(slide, "포함 케이스", BODY_X, Cm(9.2), BODY_W, Cm(0.7))
            add_l2_bullets(slide, spec["bullets"][:7], BODY_X, Cm(10.0), BODY_W, Cm(4.0))
        return

    if "table" in spec:
        headers, rows = spec["table"]
        height = Cm(9.5 if len(rows) >= 6 else 7.5)
        add_table(slide, headers, rows, BODY_X, top + Cm(0.35), BODY_W, height)
        if spec.get("bullets"):
            add_l1(slide, "설계 방향", BODY_X, top + height + Cm(0.8), BODY_W, Cm(0.6))
            add_l2_bullets(slide, spec["bullets"][:4], BODY_X, top + height + Cm(1.55), BODY_W, Cm(2.4))
        if spec.get("callout"):
            add_color_label(slide, spec["callout"], BODY_X, top + height + Cm(1.0), Cm(15), Cm(0.8), fill=RED)
        if spec.get("source"):
            add_text(slide, spec["source"], BODY_X, top + height + Cm(1.0), BODY_W, Cm(0.7), size=12, color=GRAY25)
        return

    if "columns" in spec:
        col_w = (BODY_W - Cm(1.0)) / 2
        for i, (head, items) in enumerate(spec["columns"]):
            x = BODY_X + i * (col_w + Cm(1.0))
            add_panel(slide, x, top + Cm(0.35), col_w, Cm(10.0), fill=WHITE, line=GRAY25)
            add_color_label(slide, head, x + Cm(0.35), top + Cm(0.75), Cm(5.2), Cm(0.8), fill=RED if i == 0 else ORANGE)
            add_l2_bullets(slide, items, x + Cm(0.5), top + Cm(2.0), col_w - Cm(1.0), Cm(7.5))
        return

    if "quote" in spec and "bullets" in spec:
        add_panel(slide, BODY_X, top + Cm(0.35), Cm(13.2), Cm(5.5), fill=LIGHT_GRAY, line=GRAY25)
        add_text(slide, spec["quote"], BODY_X + Cm(0.55), top + Cm(0.9), Cm(12.0), Cm(4.2), size=16, bold=True, color=BLACK)
        add_l1(slide, "주요 내용", Cm(15.0), top + Cm(0.4), Cm(16), Cm(0.7))
        add_l2_bullets(slide, spec["bullets"][:7], Cm(15.0), top + Cm(1.25), Cm(16.0), Cm(7.5))
        return

    if "code" in spec:
        add_code(slide, spec["code"], BODY_X, top + Cm(0.35), Cm(15.0), Cm(7.2))
        if spec.get("table"):
            headers, rows = spec["table"]
            add_table(slide, headers, rows, Cm(16.6), top + Cm(0.35), Cm(15.0), Cm(7.2))
        if spec.get("bullets"):
            add_l1(slide, "해석", Cm(16.6), top + Cm(0.45), Cm(15.0), Cm(0.7))
            add_l2_bullets(slide, spec["bullets"][:6], Cm(16.6), top + Cm(1.3), Cm(15.0), Cm(6.5))
        if spec.get("quote"):
            add_panel(slide, Cm(16.6), top + Cm(0.35), Cm(15.0), Cm(5.0), fill=LIGHT_GRAY, line=GRAY25)
            add_text(slide, spec["quote"], Cm(17.0), top + Cm(0.85), Cm(14.0), Cm(4.0), size=16, bold=True, color=BLACK)
        return

    if spec.get("bullets"):
        add_l1(slide, "핵심 내용", BODY_X, top + Cm(0.35), BODY_W, Cm(0.7))
        add_l2_bullets(slide, spec["bullets"], BODY_X, top + Cm(1.25), BODY_W, Cm(10.0))
        if spec.get("source"):
            add_text(slide, spec["source"], BODY_X, Cm(17.2), BODY_W, Cm(0.5), size=12, color=GRAY25)
        return


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_title_slide(prs)
    for idx, spec in enumerate(CONTENT_SLIDES, start=2):
        add_content_slide(prs, spec, idx)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
